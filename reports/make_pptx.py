from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x11, 0x17)
SURFACE = RGBColor(0x1A, 0x1D, 0x27)
BORDER  = RGBColor(0x2A, 0x2D, 0x3E)
ACCENT  = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2 = RGBColor(0x00, 0xC9, 0xA7)
ACCENT3 = RGBColor(0xFF, 0x6B, 0x6B)
TEXT    = RGBColor(0xE8, 0xEA, 0xF0)
MUTED   = RGBColor(0x8B, 0x8F, 0xA8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW  = RGBColor(0xF5, 0x9E, 0x0B)
DARK_ROW = RGBColor(0x1F, 0x22, 0x32)

W = Inches(13.33)
H = Inches(7.5)
TOTAL = 7

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def tb(slide, x, y, w, h, text, size=Pt(11), bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def label(slide, x, y, text):
    tb(slide, x, y, Inches(9), Pt(14), text, size=Pt(9), bold=True, color=ACCENT)


def slide_num(slide, n):
    tb(slide, W - Inches(1.2), Pt(10), Inches(1.1), Pt(20),
       f"{n} / {TOTAL}", size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)


def heading(slide, text, y=Inches(0.72)):
    tb(slide, Inches(0.6), y, Inches(11), Inches(0.5),
       text, size=Pt(22), bold=True, color=TEXT)


def accent_bar(slide, x, y, w=Inches(0.5)):
    add_rect(slide, x, y, w, Pt(3), fill=ACCENT)


def card(slide, x, y, w, h, title=None, title_color=ACCENT2, border=BORDER):
    add_rect(slide, x, y, w, h, fill=BG, line_color=border)
    if title:
        tb(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Pt(16),
           title.upper(), size=Pt(8), bold=True, color=title_color)
    return y + Pt(22)


def bullet_line(slide, x, y, w, text, color=MUTED, size=Pt(10.5)):
    tb(slide, x + Pt(10), y, Pt(12), Pt(16), "▸", size=Pt(9), color=ACCENT)
    tb(slide, x + Pt(24), y, w - Pt(30), Pt(26), text, size=size, color=color)
    return y + Pt(22)


def mixed_run(slide, x, y, w, h, parts, base_size=Pt(10.5)):
    """parts = list of (text, bold, color)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        r.font.bold = bold
        r.font.size = base_size
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return txb


def bar_row(slide, x, y, w, label_text, value, max_val,
            fill=ACCENT, label_w=Inches(1.1), val_fmt=None):
    track_w = w - label_w - Inches(0.55)
    bar_h   = Pt(14)
    tb(slide, x, y, label_w, bar_h, label_text,
       size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, x + label_w + Pt(6), y + Pt(2), track_w, bar_h - Pt(4),
             fill=DARK_ROW)
    fill_w = max(4, int(track_w * value / max_val))
    add_rect(slide, x + label_w + Pt(6), y + Pt(2), fill_w, bar_h - Pt(4),
             fill=fill)
    val_str = val_fmt if val_fmt else f"{value:,}"
    tb(slide, x + label_w + track_w + Pt(10), y, Inches(0.5), bar_h,
       val_str, size=Pt(9), color=TEXT)
    return y + Pt(22)


def step_circle(slide, x, y, num):
    circ = slide.shapes.add_shape(9, x, y, Pt(18), Pt(18))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num); r.font.size = Pt(8); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Calibri"


def num_circle(slide, x, y, num):
    circ = slide.shapes.add_shape(9, x, y, Pt(20), Pt(20))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num); r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Calibri"


def table_row(slide, x, y, w, cols, widths, bold=False, color=TEXT,
              bg=None, size=Pt(10.5)):
    if bg:
        add_rect(slide, x, y, w, Pt(20), fill=bg)
    cx = x + Inches(0.12)
    for text, cw in zip(cols, widths):
        tb(slide, cx, y + Pt(2), cw, Pt(18), text,
           size=size, bold=bold, color=color)
        cx += cw
    add_rect(slide, x, y + Pt(20), w, Pt(1), fill=DARK_ROW)
    return y + Pt(22)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_rect(s1, 0, 0, W, H, fill=BG)
add_rect(s1, Inches(7), 0, Inches(6.33), Inches(4), fill=RGBColor(0x10, 0x0E, 0x22))
slide_num(s1, 1)
label(s1, Inches(0.7), Inches(1.6), "NLP & SOCIAL MEDIA ANALYSIS")

tb(s1, Inches(0.7), Inches(2.05), Inches(8.5), Inches(1.4),
   'Detecting "We vs. Them" Language\non Social Media',
   size=Pt(30), bold=True)

accent_bar(s1, Inches(0.7), Inches(3.55))

tb(s1, Inches(0.7), Inches(3.75), Inches(8.5), Inches(0.4),
   "A Multi-Platform NLP Pipeline for Othering Detection",
   size=Pt(14), color=TEXT)
tb(s1, Inches(0.7), Inches(4.15), Inches(8.5), Inches(0.35),
   "Applied to Immigration Discourse across Reddit, Twitter, TikTok & Instagram",
   size=Pt(12), color=MUTED)
tb(s1, Inches(0.7), Inches(4.6), Inches(8), Inches(0.35),
   "Julian Ray-Constanty  ·  Robert Gordon University  ·  2026",
   size=Pt(11), color=MUTED)

tags = [
    ("Othering Detection",    ACCENT),
    ("26,465 posts · 4 platforms", ACCENT2),
    ("NLP Pipeline",          MUTED),
    ("Social Identity Theory", MUTED),
    ("LinearSVC · F1=0.990",  ACCENT),
    ("Temporal event study",   ACCENT3),
]
tx = Inches(0.7)
for tag_text, tag_col in tags:
    w_tag = Inches(1.65)
    add_rect(s1, tx, Inches(5.3), w_tag, Pt(22), fill=BG, line_color=tag_col)
    tb(s1, tx + Pt(4), Inches(5.32), w_tag - Pt(8), Pt(20),
       tag_text, size=Pt(8.5), color=tag_col, align=PP_ALIGN.CENTER)
    tx += w_tag + Pt(6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Aim & Objectives
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, W, H, fill=SURFACE)
slide_num(s2, 2)
label(s2, Inches(0.6), Inches(0.4), "RESEARCH FRAMING")
heading(s2, "Aim & Objectives")

CX = Inches(0.55); CY = Inches(1.45); CW = Inches(5.5); CH = Inches(5.6)
card(s2, CX, CY, CW, CH, "AIM")
tb(s2, CX + Inches(0.15), CY + Inches(0.45), CW - Inches(0.3), Inches(1.2),
   'Detect and measure "othering" language — the rhetorical move of framing an '
   'out-group as threatening against a cohesive "us" — in immigration-related '
   'discourse across four social media platforms, using a modular, reproducible NLP pipeline.',
   size=Pt(11), color=MUTED)
tb(s2, CX + Inches(0.15), CY + Inches(1.9), CW - Inches(0.3), Pt(16),
   "Theoretical backbone:", size=Pt(10), bold=True)
tb(s2, CX + Inches(0.15), CY + Inches(2.25), CW - Inches(0.3), Inches(1.1),
   "Social Identity Theory (Tajfel & Turner, 1979) — in-group/out-group "
   "categorisation manifests in pronoun choices, threat metaphors, and exclusionary framing.",
   size=Pt(10), color=MUTED, italic=True)

# Definition box
DY = CY + Inches(3.55)
add_rect(s2, CX, DY, CW, Inches(1.6), fill=RGBColor(0x12, 0x10, 0x1E),
         line_color=RGBColor(0x3A, 0x30, 0x60))
tb(s2, CX + Inches(0.15), DY + Inches(0.12), CW - Inches(0.3), Pt(14),
   "DEFINITION", size=Pt(8), bold=True, color=ACCENT)
tb(s2, CX + Inches(0.15), DY + Inches(0.38), CW - Inches(0.3), Inches(1.0),
   '"Othering" — constructing a threatening or alien "them" in contrast to a '
   'cohesive "us", legitimising exclusion. — Tajfel & Turner (1979); van Dijk (1993)',
   size=Pt(10), color=MUTED, italic=True)

RX = Inches(6.4); RY = Inches(1.45); RW = Inches(6.5)
objectives = [
    ("Measure prevalence",        " of othering language across Reddit, Twitter, TikTok & Instagram"),
    ("Test pronoun structure",    " (we / them markers) as predictor of othering intent"),
    ("Compare toxicity, emotion,", " and topic profiles across platforms and subreddits"),
    ("Validate supervised classifier", " (LinearSVC) that generalises beyond dictionary patterns"),
    ("Track discourse evolution", " over time (2020–2023) via event-aligned temporal analysis"),
]
oy = RY
for bold_part, rest in objectives:
    tb(s2, RX, oy, Pt(14), Pt(16), "▸", size=Pt(9), color=ACCENT)
    mixed_run(s2, RX + Pt(16), oy, RW - Pt(16), Pt(30),
              [(bold_part, True, TEXT), (rest, False, MUTED)], base_size=Pt(11))
    oy += Pt(54)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Datasets
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, W, H, fill=SURFACE)
slide_num(s3, 3)
label(s3, Inches(0.6), Inches(0.4), "METHODOLOGY — DATA")
heading(s3, "Datasets")

# ── Left: platform overview table ──────────────────────────────────────────
LX = Inches(0.55); LY = Inches(1.42); LW = Inches(6.5)
card(s3, LX, LY, LW, Inches(3.2), "PLATFORM CORPORA — 26,465 POSTS TOTAL",
     title_color=ACCENT, border=ACCENT)

# Column headers
hy = LY + Inches(0.45)
col_widths = [Inches(1.7), Inches(0.8), Inches(1.1), Inches(2.5)]
table_row(s3, LX + Inches(0.1), hy, LW - Inches(0.2),
          ["Platform", "Posts", "Period", "Source"],
          col_widths, bold=True, color=MUTED, size=Pt(8.5))

platforms = [
    ("Reddit",    "11,247", "2020–2023", "Arctic Shift API · 5 subreddits",   ACCENT3),
    ("Twitter",   "10,000", "2016–2021", "Immigration keywords",               RGBColor(0x1D, 0x9B, 0xF0)),
    ("Instagram",  "4,018", "2020–2022", "Immigration hashtags",               RGBColor(0xE1, 0x30, 0x6C)),
    ("TikTok",     "1,200", "2021–2023", "Immigration content (+ 19,382 trans.)", RGBColor(0xFF, 0x00, 0x50)),
]
ry2 = hy + Pt(22)
for plat, n, period, src, col in platforms:
    add_rect(s3, LX + Inches(0.1), ry2, LW - Inches(0.2), Pt(20),
             fill=RGBColor(0x14, 0x16, 0x20))
    cx2 = LX + Inches(0.22)
    for text, cw, tcol in [(plat, col_widths[0], col),
                            (n, col_widths[1], TEXT),
                            (period, col_widths[2], MUTED),
                            (src, col_widths[3], MUTED)]:
        tb(s3, cx2, ry2 + Pt(2), cw, Pt(18), text,
           size=Pt(10), bold=(text == plat), color=tcol)
        cx2 += cw
    add_rect(s3, LX + Inches(0.1), ry2 + Pt(20), LW - Inches(0.2), Pt(1), fill=DARK_ROW)
    ry2 += Pt(22)

tb(s3, LX + Inches(0.15), LY + Inches(2.85), LW - Inches(0.3), Pt(18),
   "All datasets classified with ML othering model (LinearSVC) + Detoxify toxicity.",
   size=Pt(9), color=MUTED, italic=True)

# ── Left bottom: Reddit subreddit bar chart ─────────────────────────────────
BR_Y = LY + Inches(3.35)
card(s3, LX, BR_Y, LW, Inches(3.6), "REDDIT — POST DISTRIBUTION BY SUBREDDIT")
by = BR_Y + Inches(0.5)
subs = [("r/politics",   2250, ACCENT), ("r/worldnews",  2250, ACCENT),
        ("r/ukpolitics", 2250, ACCENT), ("r/europe",     2250, ACCENT),
        ("r/immigration",2247, ACCENT2)]
for lbl, val, col in subs:
    by = bar_row(s3, LX + Inches(0.1), by, LW - Inches(0.2), lbl, val, 2250, fill=col)

tb(s3, LX + Inches(0.15), BR_Y + Inches(2.85), LW - Inches(0.3), Pt(18),
   "Event-targeted sampling: ±30 days around 45 immigration events (2020–2023)",
   size=Pt(9), color=MUTED, italic=True)

# ── Right: training corpus + preprocessing ─────────────────────────────────
RX3 = Inches(7.35); RY3 = Inches(1.42); RW3 = Inches(5.6)
card(s3, RX3, RY3, RW3, Inches(2.85), "TRAINING CORPUS — 134,459 POSTS")
bullet_items_tc = [
    "133,170 annotated posts — Kennedy et al. (2020), UC Berkeley Measuring Hate Speech",
    "1,289 Reddit posts (mixed political & immigration subreddits)",
    "Silver labels: rule-based othering detector used to generate training labels",
    "80/20 train/test split — 107,567 train · 26,892 test",
]
ty3 = RY3 + Inches(0.45)
for item in bullet_items_tc:
    ty3 = bullet_line(s3, RX3 + Inches(0.1), ty3, RW3 - Inches(0.2), item, size=Pt(10))

PPY3 = RY3 + Inches(3.0)
card(s3, RX3, PPY3, RW3, Inches(3.95), "PREPROCESSING")
prep_items = [
    "Lowercasing · URL & @mention removal",
    "Special-character normalisation, basic punctuation kept",
    "Posts < 20 chars removed → stored in clean_text column",
    "Balanced sampling across subreddits (≈ 2,250 each)",
    "Emotion (GoEmotions) available: Twitter, TikTok, Instagram",
    "Toxicity (Detoxify) available: all platforms",
]
py3 = PPY3 + Inches(0.45)
for item in prep_items:
    py3 = bullet_line(s3, RX3 + Inches(0.1), py3, RW3 - Inches(0.2), item, size=Pt(10))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Pipeline / Models
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, W, H, fill=SURFACE)
slide_num(s4, 4)
label(s4, Inches(0.6), Inches(0.4), "METHODOLOGY — PIPELINE")
heading(s4, "6-Step NLP Pipeline")

steps = [
    ("Pronoun tagging",
     "WE / THEM marker flags (binary + count + category: both / we_only / them_only / none). "
     "THEM list includes 'those people', 'people like them' for derogatory framing."),
    ("Toxicity scoring",
     "Detoxify original model: toxicity, severe_toxicity, identity_attack, "
     "insult, threat (floats 0–1)"),
    ("Emotion classification",
     "GoEmotions (monologg/bert-base-cased-goemotions-original, 27 categories, top-1). "
     "Available on Twitter, TikTok, Instagram."),
    ("Rule-based othering detector",
     "33 surface patterns across 4 categories: threat metaphors (invasion, swarm, horde…), "
     "moral exclusion, generalisation, replacement conspiracy (great replacement…). "
     "Used to generate silver training labels."),
    ("Supervised classifier",
     "TF-IDF (50k features, 1–2 grams) + LinearSVC; also tested Logistic Regression & "
     "sentence embeddings (all-MiniLM-L6-v2). 80/20 train/test split."),
    ("Topic modeling",
     "BERTopic (sentence embeddings + UMAP + HDBSCAN, min_topic_size=50). "
     "Topic assignments cross-referenced with toxicity, emotion & othering rates."),
]
col_w = Inches(6.1)
for i, (title, desc) in enumerate(steps):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.55) + col * (col_w + Inches(0.4))
    y = Inches(1.45) + row * Inches(1.9)
    add_rect(s4, x, y, col_w, Inches(1.75), fill=BG, line_color=BORDER)
    step_circle(s4, x + Inches(0.15), y + Inches(0.15), i + 1)
    tb(s4, x + Inches(0.55), y + Inches(0.12), col_w - Inches(0.65), Pt(18),
       title, size=Pt(11), bold=True)
    tb(s4, x + Inches(0.2), y + Inches(0.5), col_w - Inches(0.35), Inches(1.1),
       desc, size=Pt(9.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Classifier Performance
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, W, H, fill=SURFACE)
slide_num(s5, 5)
label(s5, Inches(0.6), Inches(0.4), "RESULTS — CLASSIFIER")
heading(s5, "Classifier Performance")

# ── Left: metrics ─────────────────────────────────────────────────────────
LX5 = Inches(0.55); LY5 = Inches(1.42); LW5 = Inches(6.0)

card(s5, LX5, LY5, LW5, Inches(1.85),
     "BEST MODEL — LINEARSVC + TF-IDF", border=ACCENT2)
for i, (val, lbl) in enumerate([("0.990", "F1 score"), ("0.994", "Precision"), ("0.985", "Recall")]):
    sx5 = LX5 + Inches(0.2) + i * Inches(1.9)
    tb(s5, sx5, LY5 + Inches(0.55), Inches(1.8), Pt(30), val,
       size=Pt(26), bold=True, color=ACCENT2)
    tb(s5, sx5, LY5 + Inches(1.1), Inches(1.8), Pt(16), lbl,
       size=Pt(9), color=MUTED)
tb(s5, LX5 + Inches(0.2), LY5 + Inches(1.48), LW5 - Inches(0.4), Pt(18),
   "Trained on 107,567 samples (5,692 othering) · Tested on 26,892 samples",
   size=Pt(9), color=MUTED, italic=True)

# Confusion matrix
CM_Y = LY5 + Inches(2.05)
card(s5, LX5, CM_Y, LW5, Inches(2.9), "CONFUSION MATRIX (TEST SET — 26,892 SAMPLES)")
# Header row
cm_cols = [Inches(1.9), Inches(1.9), Inches(1.9)]
cm_y = CM_Y + Inches(0.45)
table_row(s5, LX5 + Inches(0.1), cm_y, LW5 - Inches(0.2),
          ["", "Predicted: NOT other.", "Predicted: OTHERING"],
          cm_cols, bold=True, color=MUTED, size=Pt(8.5))
cm_y += Pt(22)
# TN/FP row
table_row(s5, LX5 + Inches(0.1), cm_y, LW5 - Inches(0.2),
          ["Actual: NOT othering", "25,461  ✓  TN", "8  ✗  FP"],
          cm_cols, color=TEXT, size=Pt(10))
cm_y += Pt(22)
# FN/TP row
table_row(s5, LX5 + Inches(0.1), cm_y, LW5 - Inches(0.2),
          ["Actual: OTHERING", "21  ✗  FN", "1,402  ✓  TP"],
          cm_cols, color=TEXT, size=Pt(10))
# Colour the good cells
tb(s5, LX5 + Inches(0.1) + cm_cols[0], CM_Y + Inches(0.9),
   cm_cols[1], Pt(18), "25,461  ✓  TN", size=Pt(10), bold=True, color=ACCENT2)
tb(s5, LX5 + Inches(0.1) + cm_cols[0] + cm_cols[1] + cm_cols[2], CM_Y + Inches(1.12),
   Inches(1.5), Pt(18), "1,402  ✓  TP", size=Pt(10), bold=True, color=ACCENT2)
tb(s5, LX5 + Inches(0.15) + cm_cols[0] + cm_cols[1], CM_Y + Inches(0.9),
   cm_cols[2], Pt(18), "8  ✗  FP", size=Pt(10), color=ACCENT3)
tb(s5, LX5 + Inches(0.1) + cm_cols[0], CM_Y + Inches(1.12),
   cm_cols[1], Pt(18), "21  ✗  FN", size=Pt(10), color=YELLOW)

tb(s5, LX5 + Inches(0.15), CM_Y + Inches(2.55), LW5 - Inches(0.3), Pt(18),
   "Only 29 misclassifications out of 26,892.  False positive rate: 0.03%",
   size=Pt(9), color=MUTED, italic=True)

# Model comparison
MC_Y = CM_Y + Inches(3.05)
card(s5, LX5, MC_Y, LW5, Inches(1.3), "MODEL COMPARISON")
table_row(s5, LX5 + Inches(0.1), MC_Y + Inches(0.43), LW5 - Inches(0.2),
          ["Model", "Features", "Precision", "Recall", "F1"],
          [Inches(1.5), Inches(1.3), Inches(0.9), Inches(0.9), Inches(0.9)],
          bold=True, color=MUTED, size=Pt(8.5))
table_row(s5, LX5 + Inches(0.1), MC_Y + Inches(0.43) + Pt(22), LW5 - Inches(0.2),
          ["LinearSVC ★", "TF-IDF", "0.994", "0.985", "0.990"],
          [Inches(1.5), Inches(1.3), Inches(0.9), Inches(0.9), Inches(0.9)],
          bold=True, color=ACCENT2)
table_row(s5, LX5 + Inches(0.1), MC_Y + Inches(0.43) + Pt(44), LW5 - Inches(0.2),
          ["Logistic Reg.", "TF-IDF", "0.943", "0.973", "0.958"],
          [Inches(1.5), Inches(1.3), Inches(0.9), Inches(0.9), Inches(0.9)],
          color=MUTED)

# ── Right: silver labels + topic ───────────────────────────────────────────
RX5 = Inches(6.9); RY5 = Inches(1.42); RW5 = Inches(6.0)
card(s5, RX5, RY5, RW5, Inches(2.75), "LABEL METHODOLOGY — SILVER LABELS")
sl_items = [
    ("Training labels",        " generated by 33-pattern rule-based detector, not human annotators"),
    ("High F1 (0.990)",        " reflects how well the model imitates the rules — not gold human judgement"),
    ("Rules from SIT/CDA:",    " threat metaphors, generalisation, moral exclusion, conspiracy framing"),
    ("Model generalises",      " beyond exact pattern matching — captures paraphrase & lexical variation"),
    ("Predicted othering rate:", " 5.29% (vs 5.30% rule-based) on full training corpus"),
]
sy5 = RY5 + Inches(0.45)
for bold_p, rest_p in sl_items:
    tb(s5, RX5 + Pt(10), sy5, Pt(12), Pt(16), "▸", size=Pt(9), color=ACCENT)
    mixed_run(s5, RX5 + Pt(24), sy5, RW5 - Pt(30), Pt(26),
              [(bold_p, True, TEXT), (rest_p, False, MUTED)])
    sy5 += Pt(28)

# Classifier note
NT_Y = RY5 + Inches(2.9)
card(s5, RX5, NT_Y, RW5, Inches(1.75), "TOXICITY & EMOTION (REDDIT CORPUS)")
tox_items = [
    ("Avg toxicity 0.020", " on Reddit vs 0.595 on hate speech reference corpus (30× lower)"),
    ("Dominant emotions:", " neutral (79.8%), annoyance, curiosity across full corpus"),
    ("Othering posts", " skew toward anger & disgust vs neutral baseline"),
]
ny5 = NT_Y + Inches(0.42)
for bold_p, rest_p in tox_items:
    tb(s5, RX5 + Pt(10), ny5, Pt(12), Pt(16), "▸", size=Pt(9), color=ACCENT)
    mixed_run(s5, RX5 + Pt(24), ny5, RW5 - Pt(30), Pt(26),
              [(bold_p, True, TEXT), (rest_p, False, MUTED)])
    ny5 += Pt(26)

# BERTopic card
BT_Y = NT_Y + Inches(1.9)
card(s5, RX5, BT_Y, RW5, Inches(2.0), "TOPIC MODELING — BERTOPIC")
tb(s5, RX5 + Inches(0.15), BT_Y + Inches(0.45), RW5 - Inches(0.3), Inches(1.3),
   "1,334 topics identified. Othering concentrates in clusters around racial identity, "
   "migration, antisemitism, and geopolitics — nearly absent from other thematic areas. "
   "Topics cross-referenced with othering rate, toxicity mean, and dominant emotion per cluster.",
   size=Pt(10.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Platform Findings
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, W, H, fill=SURFACE)
slide_num(s6, 6)
label(s6, Inches(0.6), Inches(0.4), "RESULTS — PLATFORMS")
heading(s6, "Cross-Platform Findings")

# ── Left top: platform othering bar chart ─────────────────────────────────
L6 = Inches(0.55); T6 = Inches(1.42); W6 = Inches(5.8)
card(s6, L6, T6, W6, Inches(2.3), "OTHERING RATE BY PLATFORM")
pb = T6 + Inches(0.48)
plat_bars = [
    ("TikTok",    3.58, ACCENT3,  "3.58%"),
    ("Instagram", 1.77, ACCENT,   "1.77%"),
    ("Twitter",   1.66, ACCENT,   "1.66%"),
    ("Reddit",    0.82, ACCENT2,  "0.82%"),
]
for plat, val, col, fmt in plat_bars:
    pb = bar_row(s6, L6 + Inches(0.1), pb, W6 - Inches(0.2),
                 plat, val, 3.58, fill=col, val_fmt=fmt)
tb(s6, L6 + Inches(0.15), T6 + Inches(1.95), W6 - Inches(0.3), Pt(18),
   "TikTok shows 4× higher othering rate than Reddit. r/immigration (0.04%) is lowest across all.",
   size=Pt(9), color=MUTED, italic=True)

# ── Left bottom: Reddit subreddit breakdown ────────────────────────────────
S6Y = T6 + Inches(2.45)
card(s6, L6, S6Y, W6, Inches(4.45), "REDDIT — OTHERING RATE BY SUBREDDIT")
sb = S6Y + Inches(0.48)
sub_bars = [
    ("r/worldnews",  1.56, ACCENT3,  "1.56%"),
    ("r/europe",     1.07, ACCENT,   "1.07%"),
    ("r/ukpolitics", 0.76, ACCENT,   "0.76%"),
    ("r/politics",   0.67, ACCENT,   "0.67%"),
    ("r/immigration",0.04, ACCENT2,  "0.04%"),
]
for lbl, val, col, fmt in sub_bars:
    sb = bar_row(s6, L6 + Inches(0.1), sb, W6 - Inches(0.2),
                 lbl, val, 1.56, fill=col, val_fmt=fmt)
tb(s6, L6 + Inches(0.15), S6Y + Inches(2.0), W6 - Inches(0.3), Inches(1.8),
   "Highest othering in general news communities (worldnews, europe) — "
   "not in immigration-specific communities. r/worldnews leads at 1.56%, "
   "r/immigration is near zero (0.04%). Suggests othering is driven by "
   "political/international framing, not by immigration communities themselves.",
   size=Pt(9.5), color=MUTED)

# ── Right top: pronoun effect ──────────────────────────────────────────────
R6 = Inches(6.7); RW6 = Inches(6.3)
card(s6, R6, T6, RW6, Inches(3.0), "PRONOUN TYPE VS. OTHERING RATE (FULL CORPUS)", border=ACCENT)
table_row(s6, R6 + Inches(0.1), T6 + Inches(0.43), RW6 - Inches(0.2),
          ["Pronoun structure", "Share of posts", "Othering rate"],
          [Inches(2.5), Inches(1.5), Inches(1.9)],
          bold=True, color=MUTED, size=Pt(8.5))
pron_rows = [
    ("both (we + them)", "6.7%",  "22.3%", ACCENT2),
    ("them only",        "18.5%", "12.3%", TEXT),
    ("we only",          "11.3%", " 6.9%", TEXT),
    ("none",             "63.4%", " 1.2%", MUTED),
]
pr_y = T6 + Inches(0.43) + Pt(22)
for plbl, share, rate, col in pron_rows:
    bold_it = (col == ACCENT2)
    table_row(s6, R6 + Inches(0.1), pr_y, RW6 - Inches(0.2),
              [plbl, share, rate],
              [Inches(2.5), Inches(1.5), Inches(1.9)],
              bold=bold_it, color=col)
    pr_y += Pt(22)
tb(s6, R6 + Inches(0.15), T6 + Inches(2.68), RW6 - Inches(0.3), Pt(20),
   "Posts with both pronouns are 19× more likely to contain othering language than pronoun-free posts.",
   size=Pt(9), color=MUTED, italic=True)

# ── Right bottom: toxicity table ──────────────────────────────────────────
TX6Y = T6 + Inches(3.15)
card(s6, R6, TX6Y, RW6, Inches(3.7), "TOXICITY PROFILE — PLATFORM COMPARISON")
table_row(s6, R6 + Inches(0.1), TX6Y + Inches(0.43), RW6 - Inches(0.2),
          ["Platform", "Avg toxicity", "Hate speech ref."],
          [Inches(2.0), Inches(1.8), Inches(2.0)],
          bold=True, color=MUTED, size=Pt(8.5))
tox_rows = [
    ("TikTok",    "0.123", ACCENT3),
    ("Instagram", "0.091", TEXT),
    ("Reddit",    "0.020", TEXT),
    ("Twitter",   "0.003", TEXT),
]
tx_y = TX6Y + Inches(0.43) + Pt(22)
for plat, tval, col in tox_rows:
    table_row(s6, R6 + Inches(0.1), tx_y, RW6 - Inches(0.2),
              [plat, tval, "—" if plat != "TikTok" else "0.595 (ref.)"],
              [Inches(2.0), Inches(1.8), Inches(2.0)], color=col)
    tx_y += Pt(22)
tb(s6, R6 + Inches(0.15), TX6Y + Inches(2.1), RW6 - Inches(0.3), Inches(1.3),
   "All platforms far below the hate speech reference corpus (avg 0.595). "
   "TikTok shows the highest toxicity (0.123), consistent with its highest othering rate. "
   "Twitter toxicity is exceptionally low (0.003) — immigration keyword-sampled corpus "
   "may over-represent policy/neutral content.",
   size=Pt(9.5), color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — References
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, W, H, fill=SURFACE)
slide_num(s7, 7)
label(s7, Inches(0.6), Inches(0.4), "BIBLIOGRAPHY")
heading(s7, "References")

refs = [
    ("Baumgartner, J., Zannettou, S., Keegan, B., Squire, M., & Blackburn, J. (2020).",
     "The Pushshift Reddit dataset. Proceedings of ICWSM 2020.",
     "Background Reddit data. Primary Reddit corpus collected via Arctic Shift API (arctic-shift.photon-reddit.com)."),
    ("Grootendorst, M. (2022).",
     "BERTopic: Neural topic modeling with a class-based TF-IDF procedure. arXiv:2203.05794.",
     "Topic modeling framework — Step 6 of the pipeline."),
    ("Demszky, D. et al. (2020).",
     "GoEmotions: A dataset of fine-grained emotions. ACL 2020.",
     "27-category emotion model used in Step 3 (monologg/bert-base-cased-goemotions-original)."),
]

ry7 = Inches(1.45)
for i, (author, title, note) in enumerate(refs):
    num_circle(s7, Inches(0.55), ry7 + Pt(2), i + 1)
    txb = s7.shapes.add_textbox(Inches(1.0), ry7, Inches(12), Pt(56))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = author + " "
    r1.font.bold = True; r1.font.size = Pt(10.5); r1.font.color.rgb = TEXT; r1.font.name = "Calibri"
    r2 = p.add_run(); r2.text = title
    r2.font.size = Pt(10.5); r2.font.italic = True; r2.font.color.rgb = MUTED; r2.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r3 = p2.add_run(); r3.text = "→ " + note
    r3.font.size = Pt(9.5); r3.font.color.rgb = ACCENT2; r3.font.name = "Calibri"
    ry7 += Inches(1.1)


# ── Save ───────────────────────────────────────────────────────────────────
out = "reports/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
